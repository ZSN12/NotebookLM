"""PPT service: text extraction, keyword analysis, and slide-to-image rendering.

Slide images are saved as PNG files to disk (not base64 in the database).
"""

import base64
import logging
import os
import re
import io
import shutil
from pathlib import Path
from typing import List, Optional, Dict, Tuple
from collections import Counter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

from app.config import FONTS_DIR, QWEN_VL_API_KEY, DASHSCOPE_API_KEY

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Cross-platform slide renderer (python-pptx + Pillow)
# ---------------------------------------------------------------------------

# Known Chinese font name aliases — PPT files authored on Chinese Windows
# reference these names, but they may not exist on the deployment server.
_FONT_NAME_ALIASES: Dict[str, List[str]] = {
    "SimHei":         ["Noto Sans SC", "WenQuanYi Micro Hei", "SimHei"],
    "Microsoft YaHei": ["Noto Sans SC", "WenQuanYi Micro Hei", "Microsoft YaHei"],
    "SimSun":         ["Noto Serif SC", "Noto Sans SC", "SimSun"],
    "FangSong":       ["Noto Serif SC", "Noto Sans SC", "FangSong"],
    "KaiTi":          ["Noto Serif SC", "Noto Sans SC", "KaiTi"],
    "NSimSun":        ["Noto Serif SC", "Noto Sans SC"],
    "黑体":            ["Noto Sans SC", "SimHei"],
    "微软雅黑":         ["Noto Sans SC", "Microsoft YaHei"],
    "宋体":            ["Noto Serif SC", "SimSun"],
    "楷体":            ["Noto Serif SC", "KaiTi"],
}

# Module-level font resolution cache: (name, bold, italic) -> file path
_font_cache: Dict[Tuple[str, bool, bool], str] = {}
_system_fonts_scanned = False
_system_fonts: Dict[str, str] = {}  # lowercase name -> path


def _scan_system_fonts() -> None:
    """Walk known OS font directories and build a name→path index."""
    global _system_fonts_scanned, _system_fonts

    if _system_fonts_scanned:
        return
    _system_fonts_scanned = True

    if os.name == "nt":
        roots = [
            Path(os.environ.get("WINDIR", r"C:\Windows")) / "Fonts",
        ]
    elif os.uname().sysname == "Darwin":
        roots = [
            Path("/System/Library/Fonts"),
            Path("/Library/Fonts"),
            Path.home() / "Library/Fonts",
        ]
    else:  # Linux / other
        roots = [
            Path("/usr/share/fonts"),
            Path("/usr/local/share/fonts"),
            Path.home() / ".fonts",
        ]

    for root in roots:
        if not root.is_dir():
            continue
        for path in root.rglob("*"):
            if not path.is_file():
                continue
            suffix = path.suffix.lower()
            if suffix not in (".ttf", ".ttc", ".otf"):
                continue
            try:
                # Prefer full font name from Pillow when available
                try:
                    f = ImageFont.truetype(str(path), 12)
                    family = (f.font.family or "").lower()
                except Exception:
                    family = ""
                name = f"{family}|{path.stem.lower()}"
                _system_fonts[name] = str(path)
                # Also index by stem alone for simple lookups
                stem = path.stem.lower().replace(" ", "").replace("-", "").replace("_", "")
                if stem not in _system_fonts:
                    _system_fonts[stem] = str(path)
                # Index by full path stem for direct matches
                _system_fonts[path.stem.lower()] = str(path)
            except Exception:
                pass


def _resolve_font_path(font_name: str, bold: bool = False, italic: bool = False) -> Optional[str]:
    """Resolve a font name to an available .ttf/.otf path.

    Searches in order:
      1. Cached result keyed by (name, bold, italic)
      2. Bundled font in FONTS_DIR (Noto Sans SC for CJK)
      3. System fonts found by _scan_system_fonts()
      4. Alias expansion (SimHei → Noto Sans SC, etc.)
    """
    cache_key = (font_name or "", bold, italic)
    if cache_key in _font_cache:
        return _font_cache[cache_key]

    result = _resolve_font_path_inner(font_name, bold, italic)
    _font_cache[cache_key] = result or ""
    return result


def _resolve_font_path_inner(font_name: str, bold: bool, italic: bool) -> Optional[str]:
    """Inner resolver without caching."""
    name_lower = (font_name or "").lower().strip()

    # ---- 1. Bundled fonts ----
    bundled = {
        "noto sans sc": "NotoSansSC-Regular.ttf",
        "notosanssc": "NotoSansSC-Regular.ttf",
        "notosanssc-regular": "NotoSansSC-Regular.ttf",
        "noto sans cjk sc": "NotoSansSC-Regular.ttf",
    }
    if name_lower in bundled:
        p = FONTS_DIR / bundled[name_lower]
        if p.is_file():
            return str(p)

    # Generic fallback: any .ttf in FONTS_DIR
    for f in list(FONTS_DIR.glob("*.ttf")) + list(FONTS_DIR.glob("*.ttc")) + list(FONTS_DIR.glob("*.otf")):
        return str(f)
    for f in list(FONTS_DIR.glob("*.otf")):
        return str(f)

    # ---- 2. System fonts ----
    _scan_system_fonts()
    clean = name_lower.replace(" ", "").replace("-", "").replace("_", "")
    # Search keys that contain the cleaned name
    for key, path in _system_fonts.items():
        if clean and clean in key:
            return path
    # Search by stem
    for key, path in _system_fonts.items():
        if name_lower and name_lower in key:
            return path

    # ---- 3. Alias expansion ----
    aliases = _FONT_NAME_ALIASES.get(font_name, []) if font_name else []
    # Also try case-insensitive alias match
    if not aliases:
        for alias_key, targets in _FONT_NAME_ALIASES.items():
            if alias_key.lower() == name_lower:
                aliases = targets
                break
    for alias in aliases:
        p = _resolve_font_path_inner(alias, bold, italic)
        if p:
            return p

    # ---- 4. Desperate: return any font in FONTS_DIR ----
    for f in FONTS_DIR.iterdir():
        if f.suffix.lower() in (".ttf", ".otf", ".ttc"):
            return str(f)

    # ---- 5. Use Pillow default ----
    return None


def _emu_to_px(emu: int, dpi_scale: float) -> int:
    """Convert EMU (English Metric Units) to pixels at the given scale."""
    return int(round(emu * dpi_scale))


def _rgb_to_tuple(c: RGBColor | None, default=(0, 0, 0)) -> Tuple[int, int, int]:
    """Convert python-pptx RGBColor to a PIL-friendly (R, G, B) tuple."""
    if c is None:
        return default
    try:
        return (int(c[0]), int(c[1]), int(c[2]))
    except Exception:
        return default


def _render_background(slide, draw: ImageDraw.Draw, canvas: Image.Image) -> None:
    """Paint the slide background onto the canvas."""
    from pptx.oxml.ns import qn

    try:
        bg = slide.background
        fill_elem = bg.fill
        # Try to read solid fill
        try:
            srgb = fill_elem.fore_color.rgb if fill_elem.fore_color else None
        except AttributeError:
            srgb = None
        if srgb:
            color = _rgb_to_tuple(srgb)
            draw.rectangle([(0, 0), canvas.size], fill=color)
            return
    except Exception:
        pass

    # Default: white background
    draw.rectangle([(0, 0), canvas.size], fill=(255, 255, 255))


def _load_font(font_name: str, size_pt: float, bold: bool = False, italic: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    """Load a PIL font, falling back to default if unavailable."""
    size_px = int(round(size_pt * 1.333))  # pt → px approx
    path = _resolve_font_path(font_name, bold, italic)
    if not path:
        return ImageFont.load_default()

    # .ttc files need an explicit index (PIL fails silently without one)
    if path.lower().endswith('.ttc'):
        for idx in range(4):  # try first 4 faces
            try:
                return ImageFont.truetype(path, size_px, index=idx)
            except Exception:
                continue

    try:
        return ImageFont.truetype(path, size_px)
    except OSError:
        # OTF with CFF outlines can fail in some PIL builds — try default
        pass
    except Exception:
        pass

    return ImageFont.load_default()


def _clamp_color(rgb_tuple: Tuple[int, int, int]) -> Tuple[int, int, int]:
    """Ensure all channels are in [0, 255]."""
    return (
        max(0, min(255, rgb_tuple[0])),
        max(0, min(255, rgb_tuple[1])),
        max(0, min(255, rgb_tuple[2])),
    )


def _render_text_box(draw: ImageDraw.Draw, shape, canvas: Image.Image, dpi_scale: float) -> None:
    """Render the text content of a shape that has a text frame."""
    if not shape.has_text_frame:
        return

    tf = shape.text_frame
    x = _emu_to_px(shape.left, dpi_scale)
    y = _emu_to_px(shape.top, dpi_scale)
    w = _emu_to_px(shape.width, dpi_scale)
    h = _emu_to_px(shape.height, dpi_scale)

    # Clamp to canvas (when a canvas is supplied — auto_shapes may pass None)
    if canvas is not None:
        if x >= canvas.width or y >= canvas.height:
            return
        max_x = min(x + w, canvas.width)
        max_y = min(y + h, canvas.height)
        w = max_x - x
        h = max_y - y
        if w <= 0 or h <= 0:
            return

    y_cursor = y
    default_font_size_pt = 14

    for para in tf.paragraphs:
        if y_cursor >= y + h:
            break

        align = para.alignment  # PP_ALIGN enum
        para_text = para.text or ""
        if not para_text.strip():
            y_cursor += _emu_to_px(Pt(default_font_size_pt + 4), dpi_scale)  # vertical gap
            continue

        # Build runs with formatting data
        runs_data = []
        for run in para.runs:
            text = run.text or ""
            if not text:
                continue
            font_name = run.font.name or tf.paragraphs[0].runs[0].font.name if tf.paragraphs and tf.paragraphs[0].runs else None
            font_size_pt = run.font.size
            if font_size_pt is not None:
                font_size_pt = font_size_pt / 12700.0  # EMU to pt
            else:
                font_size_pt = default_font_size_pt

            try:
                rgb_val = run.font.color.rgb if run.font.color else None
            except AttributeError:
                rgb_val = None
            color = _rgb_to_tuple(rgb_val)
            is_bold = run.font.bold or False
            is_italic = run.font.italic or False
            runs_data.append((text, font_name, font_size_pt, color, is_bold, is_italic))

        if not runs_data:
            y_cursor += _emu_to_px(Pt(default_font_size_pt + 2), dpi_scale)
            continue

        # Draw each run horizontally on the same baseline
        x_cursor = x
        line_height = 0
        line_parts: List[Tuple[ImageFont.FreeTypeFont | ImageFont.ImageFont, str, Tuple[int, int, int], float, float]] = []

        for text, font_name, size_pt, color, is_bold, is_italic in runs_data:
            font = _load_font(font_name or "", size_pt, is_bold, is_italic)
            line_parts.append((font, text, color, x_cursor, y_cursor))
            # Advance x
            try:
                bbox = draw.textbbox((0, 0), text, font=font)
                text_w = bbox[2] - bbox[0]
                line_h = bbox[3] - bbox[1]
            except Exception:
                text_w = len(text) * int(size_pt * 0.75)
                line_h = int(size_pt * 1.5)
            x_cursor += text_w
            line_height = max(line_height, line_h)

        # Adjust horizontal alignment
        total_text_width = x_cursor - x
        if align == 2:  # CENTER
            offset_x = (w - total_text_width) / 2
        elif align == 3:  # RIGHT
            offset_x = w - total_text_width
        else:  # LEFT or None
            offset_x = 0

        for font, text, color, orig_x, orig_y_cursor in line_parts:
            adjusted_x = orig_x + offset_x
            if adjusted_x + 10 >= x + w:
                continue
            draw.text((adjusted_x, orig_y_cursor), text, fill=color, font=font)

        y_cursor += line_height + 2  # small inter-line gap


def _render_image_shape(shape, canvas: Image.Image, dpi_scale: float) -> None:
    """Extract embedded image blob and paste onto canvas."""
    try:
        image = shape.image
        img_bytes = image.blob
        pil_img = Image.open(io.BytesIO(img_bytes)).convert("RGBA")
    except Exception as e:
        print(f"[PPT-RENDER] Failed to extract image: {e}")
        return

    x = _emu_to_px(shape.left, dpi_scale)
    y = _emu_to_px(shape.top, dpi_scale)
    w = _emu_to_px(shape.width, dpi_scale)
    h = _emu_to_px(shape.height, dpi_scale)

    if w <= 0 or h <= 0:
        return

    try:
        pil_img = pil_img.resize((w, h), Image.LANCZOS)
    except Exception:
        pil_img = pil_img.resize((max(1, w), max(1, h)))

    try:
        canvas.paste(pil_img, (x, y), pil_img if pil_img.mode == "RGBA" else None)
    except Exception as e:
        print(f"[PPT-RENDER] Failed to paste image at ({x},{y}): {e}")


def _render_auto_shape(draw: ImageDraw.Draw, shape, dpi_scale: float) -> None:
    """Render basic auto-shapes (rectangles, lines, ovals, etc.)."""
    from pptx.enum.shapes import MSO_SHAPE

    try:
        auto_type = shape.auto_shape_type
    except Exception:
        return

    x = _emu_to_px(shape.left, dpi_scale)
    y = _emu_to_px(shape.top, dpi_scale)
    w = _emu_to_px(shape.width, dpi_scale)
    h = _emu_to_px(shape.height, dpi_scale)

    # Fill color
    fill_color = None
    try:
        fill = shape.fill
        if fill is not None:
            try:
                if fill.fore_color and fill.fore_color.type is not None:
                    fill_color = _rgb_to_tuple(fill.fore_color.rgb)
            except AttributeError:
                pass
    except Exception:
        pass

    # Line color & width are fragile — hide all attribute errors
    outline_color = None
    outline_width = 1
    try:
        line = shape.line
        if line is not None:
            try:
                if line.color and line.color.type is not None:
                    outline_color = _rgb_to_tuple(line.color.rgb)
            except AttributeError:
                pass
            try:
                if line.width is not None:
                    outline_width = max(1, int(round(line.width / 12700)))
            except (AttributeError, TypeError):
                pass
    except Exception:
        pass
        pass

    # Simple enum check — MSO_SHAPE values
    if auto_type == MSO_SHAPE.RECTANGLE or auto_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        if fill_color:
            draw.rectangle([(x, y), (x + w, y + h)], fill=fill_color)
        if outline_color:
            draw.rectangle([(x, y), (x + w, y + h)], outline=outline_color, width=outline_width)
        # If shape also has text, render it
        if shape.has_text_frame:
            _render_text_box(draw, shape, None, dpi_scale)  # canvas ref not needed for draw-only

    elif auto_type == MSO_SHAPE.OVAL:
        if fill_color:
            draw.ellipse([(x, y), (x + w, y + h)], fill=fill_color)
        if outline_color:
            draw.ellipse([(x, y), (x + w, y + h)], outline=outline_color, width=outline_width)

    elif auto_type in (MSO_SHAPE.LINE, MSO_SHAPE.STRAIGHT_CONNECTOR):
        # Lines: width/height is actually endpoint in relative terms
        x2 = x + w
        y2 = y + h
        draw.line([(x, y), (x2, y2)], fill=outline_color or (0, 0, 0), width=outline_width)

    elif auto_type == MSO_SHAPE.ISOSCELES_TRIANGLE or auto_type == MSO_SHAPE.RIGHT_TRIANGLE:
        if auto_type == MSO_SHAPE.RIGHT_TRIANGLE:
            pts = [(x, y), (x, y + h), (x + w, y + h)]
        else:
            pts = [(x + w // 2, y), (x, y + h), (x + w, y + h)]
        if fill_color:
            draw.polygon(pts, fill=fill_color)
        if outline_color:
            draw.polygon(pts, outline=outline_color)

    else:
        # Minimal fallback for unrecognized shapes: draw bounding box
        if fill_color:
            draw.rectangle([(x, y), (x + w, y + h)], fill=fill_color)
        if outline_color:
            draw.rectangle([(x, y), (x + w, y + h)], outline=outline_color, width=outline_width)


def _render_table(draw: ImageDraw.Draw, shape, canvas: Image.Image, dpi_scale: float) -> None:
    """Render a table with grid lines and cell text."""
    if not shape.has_table:
        return
    table = shape.table

    x = _emu_to_px(shape.left, dpi_scale)
    y = _emu_to_px(shape.top, dpi_scale)

    rows = len(table.rows)
    cols = len(table.columns)

    # Calculate column widths and row heights from EMU
    col_widths = []
    for col in table.columns:
        col_widths.append(_emu_to_px(col.width, dpi_scale))
    row_heights = []
    for row in table.rows:
        row_heights.append(_emu_to_px(row.height, dpi_scale))

    y_cursor = y
    row_font = _load_font("", 11, False, False)

    for ri in range(rows):
        if ri >= len(row_heights):
            break
        rh = row_heights[ri]
        x_cursor = x

        # Alternate row background
        bg = (240, 240, 240) if ri % 2 == 0 else (255, 255, 255)
        draw.rectangle([(x, y_cursor), (x + sum(col_widths), y_cursor + rh)], fill=bg)

        for ci in range(cols):
            if ci >= len(col_widths):
                break
            cw = col_widths[ci]
            cell = table.cell(ri, ci)
            text = (cell.text or "").strip()

            # Draw cell border
            draw.rectangle([(x_cursor, y_cursor), (x_cursor + cw, y_cursor + rh)],
                           outline=(200, 200, 200), width=1)

            if text:
                try:
                    bbox = draw.textbbox((0, 0), text, font=row_font)
                    tw = bbox[2] - bbox[0]
                    th = bbox[3] - bbox[1]
                except Exception:
                    tw = len(text) * 7
                    th = 14
                tx = x_cursor + 4
                ty = y_cursor + (rh - th) // 2
                if tw > cw - 8:
                    # Truncate with ellipsis
                    while text and tw > cw - 8:
                        text = text[:-1]
                        try:
                            bbox = draw.textbbox((0, 0), text + "...", font=row_font)
                            tw = bbox[2] - bbox[0]
                        except Exception:
                            tw = len(text) * 7
                    text = text + "..." if text else ""
                if text and tx + tw < x_cursor + cw - 2:
                    draw.text((tx, ty), text, fill=(0, 0, 0), font=row_font)

            x_cursor += cw
        y_cursor += rh


def _render_group(draw: ImageDraw.Draw, shape, canvas: Image.Image, dpi_scale: float) -> None:
    """Recursively render shapes inside a group shape."""
    try:
        for child in shape.shapes:
            _render_shape(draw, child, canvas, dpi_scale)
    except Exception as e:
        print(f"[PPT-RENDER] Group render failed: {e}")


def _render_shape(draw: ImageDraw.Draw, shape, canvas: Image.Image, dpi_scale: float) -> None:
    """Dispatch a single shape to the correct renderer."""
    try:
        shape_type = shape.shape_type

        if shape_type == MSO_SHAPE_TYPE.GROUP:
            _render_group(draw, shape, canvas, dpi_scale)

        elif shape_type == MSO_SHAPE_TYPE.PICTURE:
            _render_image_shape(shape, canvas, dpi_scale)

        elif shape_type == MSO_SHAPE_TYPE.TABLE:
            _render_table(draw, shape, canvas, dpi_scale)

        elif shape_type in (MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER):
            _render_text_box(draw, shape, canvas, dpi_scale)

        elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            _render_auto_shape(draw, shape, dpi_scale)
            # Auto-shapes may also contain text
            if shape.has_text_frame:
                _render_text_box(draw, shape, canvas, dpi_scale)

        else:
            # Fallback: if shape has text, render the text at least
            if shape.has_text_frame:
                _render_text_box(draw, shape, canvas, dpi_scale)
            else:
                pass  # Silently skip unsupported types (charts, SmartArt, media, etc.)

    except NotImplementedError:
        pass  # Expected for unsupported shapes
    except Exception as e:
        print(f"[PPT-RENDER] Shape {shape.shape_type} render error: {e}")


def _render_slides_via_pillow(ppt_path: str, output_dir: str, slide_count: int, slide_width_px: int = 1920) -> bool:
    """Render each slide as PNG using python-pptx + Pillow (cross-platform).

    Returns True if at least one slide was rendered successfully.
    """
    rendered_any = False

    try:
        prs = Presentation(ppt_path)
    except Exception as e:
        print(f"[PPT-RENDER] Failed to open presentation: {e}")
        return False

    # Determine canvas pixel dimensions from slide dimensions (EMU)
    slide_w_emu = prs.slide_width or Emu(12192000)   # default 16:9
    slide_h_emu = prs.slide_height or Emu(6858000)
    dpi_scale = slide_width_px / slide_w_emu          # pixels per EMU
    canvas_h = int(round(slide_h_emu * dpi_scale))

    # Override: respect slide_count from parameter (number of slides in presentation)
    total = min(len(prs.slides), slide_count) if slide_count else len(prs.slides)

    for i, slide in enumerate(prs.slides, start=1):
        if i > total:
            break

        try:
            # RGBA canvas — white background
            canvas = Image.new("RGBA", (slide_width_px, canvas_h), (255, 255, 255, 255))
            draw = ImageDraw.Draw(canvas)

            _render_background(slide, draw, canvas)

            for shape in slide.shapes:
                try:
                    _render_shape(draw, shape, canvas, dpi_scale)
                except Exception as e:
                    print(f"[PPT-RENDER] Slide {i} shape error: {e}")

            # Resize if wider than 1400px (matching old COM behavior)
            if slide_width_px > 1400:
                ratio = 1400 / slide_width_px
                new_h = int(round(canvas_h * ratio))
                canvas = canvas.resize((1400, new_h), Image.LANCZOS)

            img_file = os.path.join(output_dir, f"slide_{i:02d}.png")
            canvas = canvas.convert("RGB")  # drop alpha for PNG
            canvas.save(img_file, format="PNG")
            rendered_any = True

        except Exception as e:
            print(f"[PPT-RENDER] Failed to render slide {i}: {e}")
            continue  # Next slide

    return rendered_any


# ---------------------------------------------------------------------------
# Legacy COM-based renderer — removed (was _render_slides_via_powerpoint).
# _render_slides_via_pillow is the cross-platform replacement.
# ---------------------------------------------------------------------------


def extract_text_from_slide(slide) -> str:
    """Extract all text from a slide (text boxes, tables, notes)."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    texts.append(text)
        if hasattr(shape, "table") and shape.table:
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        texts.append(text)
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            texts.append(f"备注: {notes_text}")
    return "\n".join(texts)


def extract_keywords_from_ppt(ppt_path: str, course_title: str) -> List[str]:
    """Extract potential domain keywords from PPT content."""
    prs = Presentation(ppt_path)
    all_text = ""
    for slide in prs.slides:
        all_text += "\n" + extract_text_from_slide(slide)

    chinese_terms = re.findall(r'[一-鿿]{2,6}', all_text)
    english_terms = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b', all_text)

    term_counts = Counter(chinese_terms + english_terms)

    common_words = {'这个', '那个', '我们', '他们', '什么', '这样', '所以', '就是', '没有', '一个',
                    'the', 'and', 'that', 'with'}
    keywords = [
        term for term, count in term_counts.most_common(30)
        if count >= 2 and term not in common_words and len(term) >= 2
    ]

    return keywords[:15]


def parse_ppt_to_slides(ppt_path: str, output_dir: Optional[str] = None) -> List[dict]:
    """Parse PPT and return list of slides with text and image metadata.

    If output_dir is provided, slide screenshots are saved as PNG files there
    (old files in the directory are cleaned first).  The returned dicts contain
    ``image_path`` (relative path to the PNG) instead of inline base64.

    Each slide dict:
        page: int (1-based)
        title: str
        text: str
        image_path: str  (e.g. "slide_01.png")
    """
    prs = Presentation(ppt_path)
    slide_count = len(prs.slides)

    rendered = False
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
        # Clean old renders before generating new ones
        for old in os.listdir(output_dir):
            if old.startswith("slide_") and old.endswith(".png"):
                try:
                    os.remove(os.path.join(output_dir, old))
                except Exception:
                    pass
        rendered = _render_slides_via_pillow(ppt_path, output_dir, slide_count)

    def _pick_title(slide) -> str:
        """Pick the best title candidate.

        Priority:
          1. Title / Center-Title placeholders (largest font, top of slide)
          2. Placeholders with large font
          3. Non-page-number text boxes, sorted by font size desc
        """
        candidates = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # Detect title placeholder type
            is_title_placeholder = False
            try:
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    # 1 = TITLE, 3 = CENTER_TITLE, 6 = SLIDE_NUMBER (skip)
                    is_title_placeholder = ph_type in (1, 3)
                    if ph_type == 6:  # slide number
                        continue
            except Exception:
                pass

            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if not txt:
                    continue
                # Skip pure numbers / page-number-like tokens
                if re.fullmatch(r"\d{1,3}|SUMMARY|CONTENTS|目录|内容", txt):
                    continue
                # Skip code fragments like '}', '{', single symbols
                if re.fullmatch(r"[\{\}\(\)\[\]<>/\\]+|private|public|if|return", txt):
                    continue

                # Estimate font size (EMU → pt)
                font_size = 0
                try:
                    if para.runs:
                        run = para.runs[0]
                        if run.font.size:
                            font_size = run.font.size / 12700.0  # EMU to pt
                except Exception:
                    pass
                if not font_size and shape.text_frame.paragraphs[0].runs:
                    try:
                        run = shape.text_frame.paragraphs[0].runs[0]
                        if run.font.size:
                            font_size = run.font.size / 12700.0
                    except Exception:
                        pass

                # Vertical position bonus (higher on slide = more likely title)
                top_y = 0
                try:
                    top_y = shape.top
                except Exception:
                    pass

                candidates.append((
                    is_title_placeholder,
                    font_size,
                    -top_y,   # higher (smaller y) first
                    len(txt),
                    txt,
                ))

        if not candidates:
            return ""

        # Sort: title placeholder > large font > higher position > shorter text
        candidates.sort(key=lambda x: (not x[0], -x[1], -x[2], x[3]))
        return candidates[0][4]

    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = _pick_title(slide)
        text = extract_text_from_slide(slide)

        image_path = ""
        if rendered:
            image_path = f"slide_{idx:02d}.png"

        slides.append({
            "page": idx,
            "title": title,
            "text": text,
            "image_path": image_path,
        })

    # ── Optional: enrich sparse slides with VL image description ──
    # If a slide has very little text (< 60 chars) but a rendered image,
    # ask Qwen-VL to describe the image so the aligner has more signal.
    if rendered and output_dir:
        _enrich_slides_with_vl(slides, str(output_dir))

    return slides


def _describe_slide_image(image_path: str, title: str) -> str:
    """Describe a slide image using Qwen-VL via DashScope OpenAI-compatible API.

    Returns empty string on any error so the caller can safely ignore it.
    """
    api_key = QWEN_VL_API_KEY or DASHSCOPE_API_KEY
    if not api_key:
        return ""

    try:
        # Resize to max 800px width to keep payload small
        with Image.open(image_path) as img:
            max_width = 800
            if img.width > max_width:
                ratio = max_width / img.width
                new_height = int(img.height * ratio)
                img = img.resize((max_width, new_height), Image.LANCZOS)

            buffer = io.BytesIO()
            img.save(buffer, format="JPEG", quality=85)
            img_base64 = base64.b64encode(buffer.getvalue()).decode("utf-8")

        from openai import OpenAI

        client = OpenAI(
            api_key=api_key,
            base_url="https://dashscope.aliyuncs.com/compatible-mode/v1",
        )

        prompt = (
            f"你是一位课堂助教。幻灯片标题是「{title}」。"
            "请用一句话描述这张教学幻灯片的核心内容，"
            "重点提取图中展示的概念、模型、流程或数据，"
            "不要描述排版和颜色。控制在80字以内。"
        )

        response = client.chat.completions.create(
            model="qwen-vl-plus",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{img_base64}"
                            },
                        },
                    ],
                }
            ],
            max_tokens=200,
            timeout=30,
        )

        desc = response.choices[0].message.content.strip()
        desc = desc.strip('"').strip("'").strip()
        return desc

    except Exception as e:
        logger.warning("vl_describe_failed path=%s error=%s", image_path, e)
        return ""


def _is_likely_chapter_header(slide: dict) -> bool:
    """Detect chapter/header slides — VL description is wasted on these.

    Chapter slides usually have a big title + a small number like '01',
    'CONTENTS', 'SUMMARY', etc. The rendered image is just a background
    with text overlay, so VL adds little value for alignment.
    """
    text = slide.get("text", "") or ""
    title = slide.get("title", "") or ""
    if not title:
        return False

    # Remove title and whitespace from text
    remainder = text.replace(title, "").strip()
    cleaned = re.sub(r"[\s\n\r\t]", "", remainder)

    # If remainder is empty or very short → likely chapter
    if len(cleaned) <= 8:
        return True

    # If remainder is just numbering/catalog words → chapter
    if re.fullmatch(r"(CONTENTS|SUMMARY|目录|内容|THANKYOU|\d{1,3})+", cleaned, re.I):
        return True

    return False


def _enrich_slides_with_vl(slides: List[dict], output_dir: str) -> None:
    """Enrich text-sparse slides with VL-generated descriptions."""
    for slide in slides:
        text = slide.get("text", "") or ""
        image_path = slide.get("image_path", "")
        if not image_path:
            continue

        # Only describe slides that are text-poor (< 60 chars)
        if len(text) >= 60:
            continue

        # Skip chapter/header slides to save API cost
        if _is_likely_chapter_header(slide):
            continue

        img_full = os.path.join(output_dir, image_path)
        if not os.path.exists(img_full):
            continue

        desc = _describe_slide_image(img_full, slide.get("title", ""))
        if desc:
            slide["text"] = f"{text}\n[图片描述]{desc}".strip()
            logger.info(
                "vl_enrich page=%s title=%s desc=%s",
                slide["page"],
                slide.get("title", "")[:20],
                desc[:40],
            )
